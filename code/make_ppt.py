"""直接修改参考PPTX的内容，保留完整样式"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REF = 'C:/Users/17296/Desktop/mm/music/Video-R1-Flow-modified.pptx'
OUTPUT = 'C:/Users/17296/Desktop/mm/music/答辩PPT.pptx'
FIG_DIR = os.path.join(os.path.dirname(OUTPUT), 'report_clk', 'figures')

prs = Presentation(REF)

def replace_text(slide, old_contains, new_text, size=None, bold=None, color=None, align=None):
    """替换slide中包含old_contains的文本框内容"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if old_contains in t or t == old_contains:
                p = shape.text_frame.paragraphs[0]
                # 只清第一个段落
                for run in p.runs:
                    run.text = ''
                p.text = new_text
                if size: p.font.size = Pt(size)
                if bold is not None: p.font.bold = bold
                if color: p.font.color.rgb = color
                if align: p.alignment = align
                return True
    return False

def add_img(slide, img_name, left_in, top_in, width_in):
    p = os.path.join(FIG_DIR, img_name)
    if os.path.exists(p):
        try:
            slide.shapes.add_picture(p, Inches(left_in), Inches(top_in), Inches(width_in))
        except:
            pass

def find_all_texts(slide):
    """返回slide中所有文本内容"""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append((shape, t))
    return texts

# ═══════════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════════
sl = prs.slides[0]
replace_text(sl, 'Video-R1-Flow', "基于稳健 Hausdorff 距离与\n多特征融合的音乐旋律线几何建模",
             size=38, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
replace_text(sl, '汇报人', "数学建模课程大作业")
replace_text(sl, '汇报日期', "2025-2026 学年")

# ═══════════════════════════════════════════════
# Slide 2: 目录
# ═══════════════════════════════════════════════
sl = prs.slides[1]
toc_map = {
    '背景与动机': '数据处理',
    '系统架构': 'Hausdorff 距离',
    '核心方法': '模型与实验',
    '训练与实验': '实验结果',
    '总结与展望': '结论',
}
toc_desc = {
    '数据处理': 'MIDI 文件格式 · Skyline 旋律提取 · 样本去重',
    'Hausdorff 距离': '数学定义 · max-HD / Q95-HD / MHD · 离群点问题',
    '模型与实验': '几何模型 · 描述符模型 · 嵌套验证 · 统计检验',
    '实验结果': '模型对比 · 混淆矩阵 · 关键发现',
    '结论': '有效，但有限 · Hausdorff 的定位',
}
all_texts = find_all_texts(sl)
for shape, t in all_texts:
    if t in toc_map:
        replace_text(sl, t, toc_map[t], size=18, color=RGBColor(0x0D,0x0D,0x0D))

# ═══════════════════════════════════════════════
# Slide 3: Section 01 数据处理
# ═══════════════════════════════════════════════
sl = prs.slides[2]
replace_text(sl, '背景与动机', '数据处理', size=28, color=RGBColor(0x0D,0x0D,0x0D))
# 保持 01 不变

# ═══════════════════════════════════════════════
# Slide 4: MIDI 格式 + Skyline
# ═══════════════════════════════════════════════
sl = prs.slides[3]
replace_text(sl, '视频AI的核心挑战', 'MIDI 文件格式与旋律提取',
             size=32, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
# 替换两个正文文本框
all_texts = find_all_texts(sl)
bodies = [s for s, t in all_texts if len(t) > 40]
for i, (shape, t) in enumerate(all_texts):
    if len(t) > 30 and 'MIDI' not in t:
        if i == 0:  # 第一个大段正文
            shape.text_frame.paragraphs[0].text = (
                "MIDI 文件存储的是演奏指令而非音频波形。\n"
                "核心事件为 Note On（按键），携带音高 pitch（0-127，每差 1 为一个半音）\n"
                "和力度 velocity（0-127，越重越响）两个参数。"
            )
        elif '挑战' not in t and len(t) > 30:  # 第二个正文
            shape.text_frame.paragraphs[0].text = (
                "Skyline 提取：同一时刻有多个音符时（和弦），只保留音高最高的那个作为主旋律近似。\n"
                "时间重采样：首末音对齐到 [0,1]，均匀取 N 个点。\n"
                "三维旋律点：时间 t、音高（减中位数→移调不变）、力度（加权 wv）。"
            )

# 加图片
add_img(sl, 'example_curves_cn.png', 0.8, 5.2, 11.5)

# ═══════════════════════════════════════════════
# Slide 5: 防泄漏去重
# ═══════════════════════════════════════════════
sl = prs.slides[4]
replace_text(sl, '看-思-答', '样本选择：防数据泄漏',
             size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
# 替换正文
for shape, t in find_all_texts(sl):
    if len(t) > 50:
        shape.text_frame.paragraphs[0].text = (
            "问题：数据集中同一首歌有多个版本（不同文件名、不同编曲），\n"
            "若分别进入训练集和测试集，会虚高分类精度。\n\n"
            "三重指纹去重：① 规范化曲名（去 live/remaster 后缀）\n"
            "② SHA-256 字节哈希 ③ 旋律指纹（移调不变的 skyline 哈希）\n\n"
            "跨曲风冲突组整组剔除，排除 75 组 187 文件。\n"
            "最终：五类各 100 首，共 500 首独立作品。"
        )
    elif '缺失' in t or '关键' in t:
        shape.text_frame.paragraphs[0].text = "同一首歌的不同版本会虚高分类精度"

add_img(sl, 'data_funnel_cn.png', 6.5, 1.5, 6)

# ═══════════════════════════════════════════════
# Slide 6: Section 02 Hausdorff
# ═══════════════════════════════════════════════
sl = prs.slides[5]
replace_text(sl, '系统架构', 'Hausdorff 距离及其变体', size=28, color=RGBColor(0x0D,0x0D,0x0D))

# ═══════════════════════════════════════════════
# Slide 7: Hausdorff 定义
# ═══════════════════════════════════════════════
sl = prs.slides[6]
replace_text(sl, '四层架构', 'Hausdorff 距离定义与变体', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
# 替换四个正文块
texts_7 = [t for _, t in find_all_texts(sl) if len(t) > 20]
new_bodies_7 = [
    "有向 Hausdorff：h(A,B) = max min ||a-b||\nA 的每个点找 B 里最近的点 → 取最大值",
    "双向 Hausdorff：H(A,B) = max{ h(A,B), h(B,A) }\n两个方向各算一次 → 值越大 = 曲线越不相似",
    "max-HD：取最近距离的最大值 → 最坏点主导\nQ95-HD：取 95% 分位数 → 去掉极端 5%\nMHD：取平均值 → 最稳健，论文主模型",
    "KD-tree 加速：O(mn) → O(m log n)\n同一首歌约 48 个采样点 vs 500 首 → 需高效计算",
]
idx = 0
for shape, t in find_all_texts(sl):
    if len(t) > 20 and idx < len(new_bodies_7):
        shape.text_frame.paragraphs[0].text = new_bodies_7[idx]
        idx += 1

# ═══════════════════════════════════════════════
# Slide 8: Section 03 模型与实验
# ═══════════════════════════════════════════════
sl = prs.slides[7]
replace_text(sl, '核心方法', '模型与实验', size=28, color=RGBColor(0x0D,0x0D,0x0D))
# 注意：Slide 8 实际上是第 9 张，索引 8。检查一下：
# Slide 1=idx0, 2=idx1, 3=idx2, 4=idx3, 5=idx4, 6=idx5, 7=idx6, 8=idx7, 9=idx8

# 实际上第8张slide（index 7）是"核心方法"的section页
# 让我重新核对：第8张是03/section，那第9张(index 8)应该是内容页

# ═══════════════════════════════════════════════
# Slide 9 (idx 8): 观察层 - 改为模型体系
# ═══════════════════════════════════════════════
sl = prs.slides[8]
replace_text(sl, '观察层', '模型体系', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
# 有 01, 02 两个卡片
for shape, t in find_all_texts(sl):
    if t == '01':
        shape.text_frame.paragraphs[0].text = '几何距离模型'
        for p in shape.text_frame.paragraphs: p.font.size = Pt(20); p.font.bold = True
    elif t == '02':
        shape.text_frame.paragraphs[0].text = '描述符模型'
        for p in shape.text_frame.paragraphs: p.font.size = Pt(20); p.font.bold = True
    elif '关键帧' in t or '停止' in t:
        shape.text_frame.paragraphs[0].text = '多参数 MHD + K-NN（主几何模型）\n相位对齐 RMSE（严格逐点对应）\n多变量 DTW（允许时间伸缩对齐）'
    elif 'Motion' in t:
        shape.text_frame.paragraphs[0].text = '随机森林 RF（非线性，最强单模型）\n多项逻辑回归（线性可分性检验）\nMHD + RF 概率融合'

# ═══════════════════════════════════════════════
# Slide 10 (idx 9): 思考层 - 改为实验设计
# ═══════════════════════════════════════════════
sl = prs.slides[9]
replace_text(sl, '思考层', '嵌套分组交叉验证', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
for shape, t in find_all_texts(sl):
    if t == '1':
        shape.text_frame.paragraphs[0].text = '外层验证'
        for p in shape.text_frame.paragraphs: p.font.size = Pt(20); p.font.bold = True
    elif t == '2':
        shape.text_frame.paragraphs[0].text = '内层选择'
        for p in shape.text_frame.paragraphs: p.font.size = Pt(20); p.font.bold = True
    elif t == '3':
        shape.text_frame.paragraphs[0].text = '统计检验'
        for p in shape.text_frame.paragraphs: p.font.size = Pt(20); p.font.bold = True
    elif '感知' in t or '推理' in t or '判断' in t:
        shape.text_frame.paragraphs[0].text = 'Stratified Group 5-Fold\n分层保持曲风比例 + 分组防止同曲跨折'
    elif '置信' in t or '阈值' in t or '工具' in t:
        shape.text_frame.paragraphs[0].text = '参数 N(36/48/96) wv(0/0.1/0.25/0.5)\nK(1/3/5/7/9) α(0.25/0.5/0.75)\n只在外训练集内以 Balanced Accuracy 选择'
    elif '定位' in t or '获取' in t or '工具' in t:
        shape.text_frame.paragraphs[0].text = 'Bootstrap 95% CI / McNemar-Holm\nPermutation Test / Pair AUC\n五组随机种子重复验证排序稳定性'

# ═══════════════════════════════════════════════
# Slide 11 (idx 10): 工具层 - 改为实验结果
# ═══════════════════════════════════════════════
sl = prs.slides[10]
replace_text(sl, '工具层', '实验结果', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
for shape, t in find_all_texts(sl):
    if '时间定位' in t:
        shape.text_frame.paragraphs[0].text = 'MHD 40.8% / RF 55.6% / 融合 54.0%'
    elif '细节' in t:
        shape.text_frame.paragraphs[0].text = 'RF 显著优于 MHD（p<0.0001）\n融合 ≈ RF（p=0.3222，不显著）'
    elif '运动' in t:
        shape.text_frame.paragraphs[0].text = '类内 MHD 0.3196 < 类间 0.3487\n差值 0.0292（p=0.0001）但 AUC=0.5645'
    elif '格式' in t or '标签' in t:
        shape.text_frame.paragraphs[0].text = '力度有用（wv=0.5）\n48 点优于 96 点\n节奏组最重要（消融下降 4.8pp）'

add_img(sl, 'model_comparison_cn.png', 0.5, 3.5, 12)

# ═══════════════════════════════════════════════
# Slide 12 (idx 11): 回答层 - 改为结论补充
# ═══════════════════════════════════════════════
sl = prs.slides[11]
replace_text(sl, '回答层', '模型对比总结', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
for shape, t in find_all_texts(sl):
    if '双维度' in t or '答案' in t:
        shape.text_frame.paragraphs[0].text = 'Hausdorff 能区分曲风但不足够'
    elif '结构化' in t or '证据' in t:
        shape.text_frame.paragraphs[0].text = 'MHD 40.8% vs 随机 20%（p<0.0001）有效但分布重叠大'

# ═══════════════════════════════════════════════
# Slide 13 (idx 12): 最终效果 - 改为 Hausdorff 定位
# ═══════════════════════════════════════════════
sl = prs.slides[12]
replace_text(sl, '收获', 'Hausdorff 的定位', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
for shape, t in find_all_texts(sl):
    if '看-思-答' in t or '循环' in t:
        shape.text_frame.paragraphs[0].text = '统计特征更强（RF 55.6%）'
    elif '模型观察' in t or '关键帧' in t or '思考' in t:
        shape.text_frame.paragraphs[0].text = (
            'Hausdorff 距离提供可解释的旋律形状相似性，\n'
            '适合作为多特征系统的辅助证据，\n'
            '而非单独追求最高分类精度。\n'
            '最终结论：有效，但有限。'
        )

# ═══════════════════════════════════════════════
# Slide 14 (idx 13): Section 04 结论
# ═══════════════════════════════════════════════
sl = prs.slides[13]
replace_text(sl, '训练与实验', '结论', size=28, color=RGBColor(0x0D,0x0D,0x0D))

# ═══════════════════════════════════════════════
# Slide 15 (idx 14): 训练详情 - 改为结论详情
# ═══════════════════════════════════════════════
sl = prs.slides[14]
replace_text(sl, '三维奖励', '结论：有效，但有限', size=28, bold=True, color=RGBColor(0x0D,0x0D,0x0D))
titles_15 = ['准确率', '格式', '一致性', '效率']
new_titles = ['Hausdorff 能区分曲风', '但区分能力有限', '统计特征远强于几何', 'Hausdorff 的定位']
new_bodies = [
    'MHD 40.8% > 随机 20%\n同类距离显著小于异类（p=0.0001）',
    'AUC=0.5645，分布高度重叠\n不能单独做高精度分类',
    'RF 55.6% >> MHD 40.8%（p<0.0001）\n节奏/力度/音级是关键描述符',
    '可解释的旋律形状相似性\n适合做几何佐证而非独立分类器',
]
# 找4个标题卡片和4个正文
shapes_15 = [(s, t) for s, t in find_all_texts(sl) if len(t) < 15 and '0' not in t and '奖' not in t and '结' not in t and '有' not in t]
body_texts = [(s, t) for s, t in find_all_texts(sl) if len(t) > 50]

for i, (shape, t) in enumerate(shapes_15[:4]):
    if i < len(new_titles):
        shape.text_frame.paragraphs[0].text = new_titles[i]
for i, (shape, t) in enumerate(body_texts[:4]):
    if i < len(new_bodies):
        shape.text_frame.paragraphs[0].text = new_bodies[i]

# ═══════════════════════════════════════════════
# Slide 16 (idx 15): 感谢页
# ═══════════════════════════════════════════════
sl = prs.slides[15]
for shape, t in find_all_texts(sl):
    if '感谢' in t:
        shape.text_frame.paragraphs[0].text = '感谢聆听'
        break

# ─── 保存 ───
prs.save(OUTPUT)
print(f"✅ 已保存: {OUTPUT}")
print(f"共 {len(prs.slides)} 页")
